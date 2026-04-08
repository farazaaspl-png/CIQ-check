import os, re, json, pandas as pd, csv, uuid, zipfile,base64
import logging
from pathlib import Path
from typing import List, Dict, Tuple
from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from copy import deepcopy
# from lxml import etree

from core.db.crud import DatabaseManager
from config import Config as cfg
# logging.getLogger().setLevel(logging.ERROR)
from core.exceptions import NoSensitiveItemFound
from core.utility import get_custom_logger
logger = get_custom_logger(__name__)

class PowerPointRedactor:
    
    # _InclusionDf = pd.read_csv(Path(r"services/process_ips/redaction/helperdata/Inclusion.csv"))
    # _EXCLUSION_ITEMS = pd.read_csv(Path(r"services/process_ips/redaction/helperdata/Exclusion.csv")).Exclude.to_list()
    
    _NSMAP = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
    }
    
    sanitize_data = []

    def __init__(self, filepath: Path, dafileid: uuid.UUID = None, debug: bool = False):
       
        self.filepath = Path(filepath) if not isinstance(filepath, Path) else filepath
        self.dafileid = dafileid
        self.filename = self.filepath.name
        logger.info(f"Processing PPT file: {self.filename} (ID: {self.dafileid})")
        
        self.ppt = Presentation(str(self.filepath))
        self.records: List[Dict] = []
        self.filecontent: str = ''
        # self.image_data: List[Dict] = []
        
        self.sensitive_info: List[Dict]=[]
        
        if debug:
            logger.setLevel(logging.DEBUG)
    
    def escape_custom(self, s: str) -> str:
        SPECIAL_CHARS = r'.^$*+?{}[]\|()#\''
        escaped = []
        for ch in s:
            if ch in SPECIAL_CHARS:
                escaped.append('\\' + ch)
            else:
                escaped.append(ch)
        return ''.join(escaped)

    def _replace_text_in_runs(self, text_frame, slide_idx, shape_idx, label, pattern, location="shape"):
        placeholder = f"<{cfg.PREFIX}{label.replace(' ','_')}>"
        safe_pattern = self.escape_custom(str(pattern))
        # Handle quote variations in the pattern
        if any(quote in safe_pattern for quote in '`\'’'):
            # Replace escaped single quote with a character class that matches both ' and `
            quote_variations = safe_pattern.replace("\\'", "[`'’]?")
        else:
            quote_variations = safe_pattern
        # Add word boundaries to avoid matching within words
        word_boundary_pattern = rf'\b{quote_variations}\b'
        # word_boundary_pattern = rf'(?<![a-zA-Z0-9_]){safe_pattern}(?![a-zA-Z0-9_])'
        # print(f"{label}: {word_boundary_pattern}")
        # print(f"=========================================")
        for para_idx, para in enumerate(text_frame.paragraphs):
            for run_idx, run in enumerate(para.runs):
                original_text = run.text if run.text is not None else ''
                if original_text == '':
                    continue
                # print(f"Original: {original_text}")
                try:
                    matches = list(re.finditer(word_boundary_pattern, run.text, re.IGNORECASE))
                    # print(f"Matches found: {len(matches)}/ {matches}")
                    
                    if not matches:
                        continue
                    len_diff = 0
                    for m in reversed(matches):
                        if m.lastindex:
                            match_str = m.group(1)
                            startIdx = m.start() + m.group().find(m.group(1))
                            endIdx = startIdx + len(m.group(1))
                        else:
                            match_str = m.group(0)
                            startIdx = m.start()
                            endIdx = m.end()
                            
                        record = {
                            "slide": slide_idx + 1,
                            "shape": shape_idx + 1,
                            "location": location,
                            "paragraph": para_idx + 1,
                            "run": run_idx + 1,
                            "start": startIdx,
                            "end": endIdx,
                            "label": label,
                            "sensitivetext": match_str,
                            "placeholder": placeholder,
                            "context": original_text[max(0, startIdx - 40):min(len(original_text), endIdx + 40)],
                            "fileid": str(self.dafileid) if self.dafileid else None
                        }
                        
                        self.records.append(record)
                        run.text = run.text[:startIdx+len_diff] + placeholder + run.text[endIdx+len_diff:]
                        #calculate difference in len of string for next iteration
                        len_diff = len_diff + len(placeholder)-(endIdx-startIdx)
                    
                except Exception as e:
                    logger.warning(f"Pattern '{label}' failed in {location}: {e}",exc_info=True)
                
                if run.text.startswith('<') and run.text.endswith('>'):
                    run.text = ''
                # print(f"modified Text: {run.text}")
        # print(f"=========================================")

    def _redact_presentation(self, label, pattern):
    
        for slide_idx, slide in enumerate(self.ppt.slides):

            for shape_idx, shape in enumerate(slide.shapes):

                if hasattr(shape, "text_frame") and shape.text_frame:
                    self._replace_text_in_runs(
                        shape.text_frame, 
                        slide_idx, 
                        shape_idx, 
                        label, 
                        pattern,
                        location="text_shape"
                    )
            
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            self._replace_text_in_runs(
                                cell.text_frame,
                                slide_idx,
                                shape_idx,
                                label,
                                pattern,
                                location=f"table_r{row_idx+1}_c{cell_idx+1}"
                            )
            
            try:
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                        self._replace_text_in_runs(
                            notes_slide.notes_text_frame,
                            slide_idx,
                            0,
                            label,
                            pattern,
                            location="notes"
                        )
            except Exception as e:
                logger.warning(f"Error redacting notes for slide {slide_idx + 1}: {e}",exc_info=True)
    
    def set_sensitiveinfo(self):
        logger.info(f"{self.dafileid}-->Fetching sensitive data from DB")
        db = DatabaseManager()
        df = db.get_vwtoberedacted(requestid=self.requestid, fuuid=self.fuuid, dafileid = self.dafileid)
        
        if df.empty:
            raise NoSensitiveItemFound()
        
        df = df[['category', 'sensitivetext']].drop_duplicates()
        df.sort_values(by="sensitivetext", key=lambda x: x.str.len(), inplace=True,ascending=False)        
        self.sensitive_data = df.to_dict(orient='records')
            
    def sanitize(self, **kwargs):    
        self.requestid = kwargs.get('requestid')
        self.fuuid = kwargs.get('fuuid')
        self.dafileid = kwargs.get('dafileid')
        
        self.set_sensitiveinfo()  

        for tr in self.sensitive_data:
            self._redact_presentation(label=tr['category'], pattern=tr['sensitivetext'])
        
        totalRedacted = len(self.records)
        return len(self.records) > 0, totalRedacted

    def save(self, outdir: str, filename:str = None):
        try:
            fname = filename if filename is not None else self.filepath.name
            OutPathFile = Path(os.path.join(outdir,fname))
            OutPathFile.parent.mkdir(parents=True, exist_ok=True)
            self.ppt.save(OutPathFile)  
            logger.info(f"Saved redacted PPTX: {OutPathFile}")

            if len(self.records)>0:
                db = DatabaseManager()
                db.insert_redacted_results(self.requestid, self.fuuid, self.dafileid, fname, self.records)
                logger.info(f"Redacted {len(self.records)} sensitive items")

                redactedItemsFilepath = os.path.join(outdir, f"{OutPathFile.stem}_redacted.csv")
                pd.DataFrame(pd.json_normalize(self.records)).to_csv(
                    redactedItemsFilepath,
                    index=False,
                    quoting=csv.QUOTE_ALL
                )
                logger.info(f"Saved redaction CSV: {redactedItemsFilepath}")
                return Path(OutPathFile),Path(redactedItemsFilepath)
            else:
                return Path(OutPathFile),None
        except Exception as e:
            logger.error(f"Saving PPTX failed: {e}", exc_info=True)
            raise e
        # finally:
        #     if getattr(self, "pptx_zip", None) and getattr(self.pptx_zip, "fp", None):
        #         self.pptx_zip.close()
    