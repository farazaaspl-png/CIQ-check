import fitz
import io
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from core.utility import get_custom_logger #,remove_control_chars


logger = get_custom_logger(__name__)

def apply_xref_mask(pdf, image_dict):
    """
    Applies PDF soft mask (SMask) to an image if present.
    Returns image bytes (PNG with alpha).
    """
    base_bytes = image_dict["image"]
    smask_xref = image_dict.get("smask")

    # No mask → return original
    if not smask_xref:
        return base_bytes

    try:
        # Load base image
        base = Image.open(io.BytesIO(base_bytes)).convert("RGBA")

        # Load mask image
        mask_dict = pdf.extract_image(smask_xref)
        mask = Image.open(io.BytesIO(mask_dict["image"])).convert("L")

        # Resize mask if needed
        if mask.size != base.size:
            mask = mask.resize(base.size, Image.BILINEAR)

        # Apply mask as alpha channel
        base.putalpha(mask)

        out = io.BytesIO()
        base.save(out, format="PNG")
        return out.getvalue()

    except Exception as e:
        # Fail safe: return original image
        print(f"[WARN] Failed to apply SMask: {e}")
        return base_bytes

# =====================================================
# COLOR & IMAGE UTILITIES
# =====================================================

def pdf_color_to_rgb(color_int):
    r = (color_int >> 16) & 255
    g = (color_int >> 8) & 255
    b = color_int & 255
    return r, g, b


def image_brightness(image_bytes):
    img = Image.open(io.BytesIO(image_bytes)).convert("L")
    return np.array(img).mean()


def detect_alpha(image_bytes):
    img = Image.open(io.BytesIO(image_bytes))
    if img.mode in ("RGBA", "LA"):
        alpha = np.array(img.split()[-1])
        return 1 - (alpha.mean() / 255.0)
    return None


def detect_page_background(page):
    pix = page.get_pixmap(matrix=fitz.Matrix(1, 1))
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

    top = np.array(img.crop((0, 0, pix.width, 20))).mean(axis=(0, 1))
    bottom = np.array(img.crop((0, pix.height - 20, pix.width, pix.height))).mean(axis=(0, 1))

    return tuple(map(int, top)), tuple(map(int, bottom))

def add_gradient_background(slide, top_rgb, bottom_rgb):
    fill = slide.background.fill
    fill.gradient()

    stops = fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = RGBColor(*top_rgb)

    stops[1].position = 1.0
    stops[1].color.rgb = RGBColor(*bottom_rgb)
    
def is_full_page(rect, page_area):
    return (rect.width * rect.height) / page_area > 0.8


# =====================================================
# PPT XML HELPERS
# =====================================================

def set_picture_transparency(picture, transparency):
    transparency = max(0.0, min(1.0, transparency))
    alpha = int((1 - transparency) * 100000)

    spPr = picture._element.spPr
    blipFill = spPr.find(qn("a:blipFill"))
    if blipFill is None:
        blipFill = OxmlElement("a:blipFill")
        spPr.append(blipFill)

    blip = blipFill.find(qn("a:blip"))
    if blip is None:
        blip = OxmlElement("a:blip")
        blipFill.insert(0, blip)

    for n in blip.findall(qn("a:alphaModFix")):
        blip.remove(n)

    mod = OxmlElement("a:alphaModFix")
    mod.set("amt", str(alpha))
    blip.append(mod)


def add_text_stroke(run, color=(0, 0, 0), width=10000):
    r, g, b = color
    ln = OxmlElement("a:ln")
    ln.set("w", str(width))

    solid = OxmlElement("a:solidFill")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
    solid.append(srgb)

    ln.append(solid)
    run._r.get_or_add_rPr().append(ln)


def add_text_shadow(run, offset=30000, blur=40000):
    """
    Adds an outer shadow to a text run.
    """
    rPr = run._r.get_or_add_rPr()

    shadow = OxmlElement("a:outerShdw")
    shadow.set("dist", str(offset))
    shadow.set("blurRad", str(blur))
    shadow.set("dir", "5400000")  # 90 degrees

    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", "000000")

    alpha = OxmlElement("a:alpha")
    alpha.set("val", "40000")  # 40% opacity

    srgb.append(alpha)
    shadow.append(srgb)
    rPr.append(shadow)

def rects_intersect(a, b):
    """Check if two rectangles intersect"""
    return not (
        a.x1 <= b.x0 or
        a.x0 >= b.x1 or
        a.y1 <= b.y0 or
        a.y0 >= b.y1
    )


def average_brightness_region(image_bytes, rect, page):
    """
    Compute average brightness of the region of the image
    that overlaps with the given rect (PDF coordinates).
    """
    img = Image.open(io.BytesIO(image_bytes)).convert("L")

    # Map PDF rect → image pixels
    img_w, img_h = img.size
    page_w, page_h = page.rect.width, page.rect.height

    x0 = int(rect.x0 / page_w * img_w)
    y0 = int(rect.y0 / page_h * img_h)
    x1 = int(rect.x1 / page_w * img_w)
    y1 = int(rect.y1 / page_h * img_h)

    # Clamp
    x0, y0 = max(0, x0), max(0, y0)
    x1, y1 = min(img_w, x1), min(img_h, y1)

    if x1 <= x0 or y1 <= y0:
        return 255  # assume light

    crop = img.crop((x0, y0, x1, y1))
    return np.array(crop).mean()

def extract_pdf_shapes(page):
    """
    Extract vector shapes with full rendering info and correct order.
    """
    shapes = []

    for d in page.get_drawings():
        rect = d.get("rect")
        seqno = d.get("seqno", 0)

        if rect is None:
            continue

        fill = d.get("fill")
        color = d.get("color")

        fill_opacity = d.get("fill_opacity", 1.0)
        stroke_opacity = d.get("stroke_opacity", 1.0)

        # Convert fill color
        fill_rgb = None
        if isinstance(fill, (list, tuple)) and len(fill) >= 3:
            fill_rgb = tuple(int(c * 255) for c in fill[:3])

        # Convert stroke color
        stroke_rgb = None
        if isinstance(color, (list, tuple)) and len(color) >= 3:
            stroke_rgb = tuple(int(c * 255) for c in color[:3])

        shapes.append({
            "rect": rect,
            "fill_rgb": fill_rgb,
            "stroke_rgb": stroke_rgb,
            "fill_opacity": fill_opacity,
            "stroke_opacity": stroke_opacity,
            "seqno": seqno
        })

    # 🔑 CRITICAL: preserve visual stacking
    shapes.sort(key=lambda s: s["seqno"])
    return shapes

def set_shape_fill_transparency(shape, transparency):
    """
    Correctly apply transparency to a PPT shape fill.
    transparency: 0.0 (opaque) → 1.0 (fully transparent)
    """
    transparency = max(0.0, min(1.0, transparency))
    alpha = int((1 - transparency) * 100000)

    spPr = shape._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return

    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is None:
        return

    # ❌ Remove WRONG alphaModFix if present
    for node in srgbClr.findall(qn("a:alphaModFix")):
        srgbClr.remove(node)

    # ✅ Remove existing alpha
    for node in srgbClr.findall(qn("a:alpha")):
        srgbClr.remove(node)

    # ✅ Correct alpha for shapes
    alpha_node = OxmlElement("a:alpha")
    alpha_node.set("val", str(alpha))
    srgbClr.append(alpha_node)



from pptx.enum.shapes import MSO_SHAPE

def add_shape_to_slide(slide, shape_info, page, prs):
    rect = shape_info["rect"]
    fill_rgb = shape_info["fill_rgb"]
    stroke_rgb = shape_info["stroke_rgb"]
    fill_opacity = shape_info["fill_opacity"]
    stroke_opacity = shape_info["stroke_opacity"]

    # Normalize PDF → PPT coordinates
    page_w, page_h = page.rect.width, page.rect.height
    slide_w, slide_h = prs.slide_width.pt, prs.slide_height.pt

    left = rect.x0 / page_w * slide_w
    top = rect.y0 / page_h * slide_h
    width = rect.width / page_w * slide_w
    height = rect.height / page_h * slide_h

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Pt(left), Pt(top), Pt(width), Pt(height)
    )

    # -------- Fill --------
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    
        # Apply transparency correctly
        set_shape_fill_transparency(shape, 1.0 - fill_opacity)
    else:
        shape.fill.background()

    # -------- Stroke --------
    if stroke_rgb and stroke_opacity > 0.2:
        shape.line.color.rgb = RGBColor(*stroke_rgb)
    else:
        shape.line.fill.background()


# =====================================================
# MAIN CONVERTER
# =====================================================

def pdf_to_pptx_final(pdf_path, pptx_path, mode="screen"):
    pdf = fitz.open(pdf_path)
    prs = Presentation()

    while prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    blank = prs.slide_layouts[6]

    for page_index, page in enumerate(pdf):
        prs.slide_width = Pt(page.rect.width)
        prs.slide_height = Pt(page.rect.height)
        slide = prs.slides.add_slide(blank)

        page_area = page.rect.width * page.rect.height

        # ---------- BACKGROUND ----------
        top_bg, bottom_bg = detect_page_background(page)
        add_gradient_background(slide, top_bg, bottom_bg)

        # ---------- SHAPES (NEW) ----------
        pdf_shapes = extract_pdf_shapes(page)
        for shp in pdf_shapes:
            add_shape_to_slide(slide, shp, page, prs)

        xref_added = []
        background_images = []  # store (rect, image_bytes)
        # ---------- IMAGES ----------
        for img in page.get_images(full=True):
            xref = img[0]
            image_dict = pdf.extract_image(xref)
            image_bytes = apply_xref_mask(pdf, image_dict)
            
            try:
                Image.open(io.BytesIO(image_bytes)).verify()
            except Exception:
                continue

            for rect in page.get_image_rects(xref):
                if xref in xref_added:
                    continue
                
                pic = slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    Pt(rect.x0), Pt(rect.y0),
                    Pt(rect.width), Pt(rect.height)
                )
                xref_added.append(xref)
                background_images.append((rect, image_bytes))

        # ---------- TEXT ----------
        page_dict = page.get_text("dict")

        for block in page_dict["blocks"]:
            if block["type"] != 0:
                continue

            for line in block["lines"]:
                x0, y0, x1, y1 = line["bbox"]

                box = slide.shapes.add_textbox(
                    Pt(x0), Pt(y0),
                    Pt(x1 - x0), Pt(y1 - y0)
                )

                tf = box.text_frame
                tf.clear()
                tf.auto_size = MSO_AUTO_SIZE.NONE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT

                for span in line["spans"]:
                    run = p.add_run()
                    run.text = span["text"]
                    run.font.size = Pt(span["size"])
                    run.font.name = span["font"]
                    run.font.bold = bool(span["flags"] & 2)
                    run.font.italic = bool(span["flags"] & 1)

                    r, g, b = pdf_color_to_rgb(span["color"])
                    # Detect background brightness
                    text_rect = fitz.Rect(x0, y0, x1, y1)
                    bg_brightness = None
                    for img_rect, img_bytes in background_images:
                        if rects_intersect(text_rect, img_rect):
                            bg_brightness = average_brightness_region(
                                img_bytes, text_rect, page
                            )
                            break
                        
                    # Decide text color
                    if bg_brightness is not None:
                        if bg_brightness < 130:   # dark background
                            run.font.color.rgb = RGBColor(255, 255, 255)
                        else:                     # light background
                            run.font.color.rgb = RGBColor(0, 0, 0)
                    else:
                        # Fallback to PDF color
                        run.font.color.rgb = RGBColor(r, g, b)

                    # Stroke for dark text
                    if r < 80 and g < 80 and b < 80:
                        add_text_stroke(run, (255, 255, 255), 8000)

                    # Shadow for headers
                    if span["size"] > 18:
                        add_text_shadow(run)

        logger.info(f"Processed page {page_index + 1}/{len(pdf)}")

    prs.save(pptx_path)
    logger.info(f"\n✅ FINAL PPTX GENERATED: {pptx_path}")


# =====================================================
# RUN
# =====================================================
# from pathlib import Path
# PDF_IN=Path(r"C:\Users\Lekhnath_Pandey\CIQ\ip_content_management\ip_content_management\input\classificationandredaction\pdf-pptx-3b7e00fd-d672-4d46-83e4-08de38a2153c\Future State Recommendations Sept 2025_V1.0.pdf")
# PPTX_OUT=Path(r"C:\Users\Lekhnath_Pandey\CIQ\ip_content_management\ip_content_management\input\classificationandredaction\pdf-pptx-3b7e00fd-d672-4d46-83e4-08de38a2153c\Future State Recommendations Sept 2025_V1.0.pptx")
# if __name__ == "__main__":
#     pdf_to_pptx_final(
#         pdf_path=PDF_IN,
#         pptx_path=PPTX_OUT,
#         mode="screen"   # change to "print" if needed
#     )


